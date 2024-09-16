import sys
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QPushButton, 
                             QTextEdit, QFileDialog, QMessageBox, QLabel, QComboBox)
from PyQt5.QtCore import Qt
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from templates import TEMPLATES

class SlideGenerator:
    def __init__(self, template):
        self.prs = Presentation()
        self.layouts = {
            'title': self.prs.slide_layouts[0],
            'content': self.prs.slide_layouts[1],
            'section': self.prs.slide_layouts[2],
        }
        self.template = template

    def parse_input(self, input_text):
        slides = []
        lines = input_text.strip().split('\n')
        
        title = lines[0][2:].strip()
        subtitle = lines[1][2:].strip() if len(lines) > 1 and lines[1].startswith('- ') else ""
        slides.append({"type": "title", "title": title, "subtitle": subtitle})
        
        current_slide = {"type": "", "title": "", "content": []}
        for line in lines[2:]:
            line = line.strip()
            if line.startswith('# '):
                if current_slide["title"]:
                    slides.append(current_slide)
                    current_slide = {"type": "", "title": "", "content": []}
                current_slide["title"] = line[2:]
                current_slide["type"] = "content"
            elif line.startswith('- ') or line.startswith('• '):
                current_slide["content"].append(line)
            elif line:
                current_slide["content"].append(line)
        if current_slide["title"]:
            slides.append(current_slide)
        return slides

    def create_slide(self, slide_data):
        slide = self.prs.slides.add_slide(self.layouts[slide_data["type"]])
        title = slide.shapes.title
        title.text = slide_data["title"]

        if slide_data["type"] == "title":
            self.format_title_slide(slide, title, slide_data["subtitle"])
        elif slide_data["type"] in ["content", "section"]:
            self.format_content_slide(slide, slide_data["content"])

        self.template.apply_style(slide)
        return slide

    def format_title_slide(self, slide, title, subtitle):
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = subtitle

    def format_content_slide(self, slide, content):
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for item in content:
            p = tf.add_paragraph()
            if item.startswith('- ') or item.startswith('• '):
                p.text = item[2:].strip()
                p.level = 0
            else:
                p.text = item
                p.level = 1
            p.alignment = self.template.content_alignment

    def generate_presentation(self, input_text, output_file):
        slides_data = self.parse_input(input_text)
        for slide_data in slides_data:
            self.create_slide(slide_data)
        self.prs.save(output_file)

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PowerPoint Slide Generator")
        self.setGeometry(100, 100, 800, 600)
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f0f0f0;
            }
            QLabel {
                font-size: 24px;
                color: #333333;
            }
            QTextEdit {
                font-size: 16px;
                border: 1px solid #cccccc;
                border-radius: 5px;
            }
            QPushButton {
                font-size: 14px;
                padding: 8px 16px;
                border-radius: 5px;
            }
            QPushButton#generate {
                background-color: #4CAF50;
                color: white;
            }
            QPushButton#generate:hover {
                background-color: #45a049;
            }
            QPushButton#clear {
                background-color: #f44336;
                color: white;
            }
            QPushButton#clear:hover {
                background-color: #da190b;
            }
            QComboBox {
                font-size: 14px;
                padding: 5px;
            }
        """)

        self.central_widget = QWidget()
        self.setCentralWidget(self.central_widget)
        self.layout = QVBoxLayout(self.central_widget)

        self.title_label = QLabel("PowerPoint Slide Generator")
        self.title_label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.title_label)

        self.text_edit = QTextEdit()
        self.layout.addWidget(self.text_edit)

        self.options_layout = QHBoxLayout()

        self.template_combo = QComboBox()
        self.template_combo.addItems(TEMPLATES.keys())
        self.options_layout.addWidget(QLabel("Template:"))
        self.options_layout.addWidget(self.template_combo)

        self.layout.addLayout(self.options_layout)

        self.button_layout = QHBoxLayout()
        self.clear_button = QPushButton("Clear Text")
        self.clear_button.setObjectName("clear")
        self.clear_button.clicked.connect(self.clear_text)
        self.button_layout.addWidget(self.clear_button)

        self.generate_button = QPushButton("Generate PowerPoint")
        self.generate_button.setObjectName("generate")
        self.generate_button.clicked.connect(self.generate_ppt)
        self.button_layout.addWidget(self.generate_button)

        self.layout.addLayout(self.button_layout)

    def clear_text(self):
        self.text_edit.clear()

    def generate_ppt(self):
        input_text = self.text_edit.toPlainText()
        
        if not input_text.strip():
            QMessageBox.critical(self, "Error", "Please enter some text for the slides.")
            return

        file_name, _ = QFileDialog.getSaveFileName(self, "Save PowerPoint Presentation", "", "PowerPoint Presentation (*.pptx)")
        if not file_name:
            return

        if not file_name.lower().endswith('.pptx'):
            file_name += '.pptx'

        template_name = self.template_combo.currentText()
        template = TEMPLATES[template_name]
        generator = SlideGenerator(template)

        try:
            generator.generate_presentation(input_text, file_name)
            QMessageBox.information(self, "Success", f"Presentation saved as {file_name}")
        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred: {str(e)}")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())