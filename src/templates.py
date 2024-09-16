from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

class Template:
    def __init__(self, name, theme_color, background_color, font_family, title_font_size, content_font_size, 
                 title_alignment=PP_ALIGN.LEFT, content_alignment=PP_ALIGN.LEFT, 
                 title_bold=False, content_bold=False, gradient=None, shape=None):
        self.name = name
        self.theme_color = theme_color
        self.background_color = background_color
        self.font_family = font_family
        self.title_font_size = title_font_size
        self.content_font_size = content_font_size
        self.title_alignment = title_alignment
        self.content_alignment = content_alignment
        self.title_bold = title_bold
        self.content_bold = content_bold
        self.gradient = gradient
        self.shape = shape

    def apply_style(self, slide):
        background = slide.background
        fill = background.fill
        
        if self.gradient:
            fill.gradient()
            fill.gradient_stops[0].color.rgb = self.gradient[0]
            fill.gradient_stops[1].color.rgb = self.gradient[1]
        else:
            fill.solid()
            fill.fore_color.rgb = self.background_color

        if self.shape:
            left = top = Inches(0)
            width = height = Inches(10)
            shape = slide.shapes.add_shape(self.shape, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme_color
            shape.line.color.rgb = self.theme_color

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    if shape == slide.shapes.title:
                        paragraph.alignment = self.title_alignment
                        font_size = self.title_font_size
                        is_bold = self.title_bold
                    else:
                        paragraph.alignment = self.content_alignment
                        font_size = self.content_font_size
                        is_bold = self.content_bold

                    for run in paragraph.runs:
                        font = run.font
                        font.name = self.font_family
                        font.size = Pt(font_size)
                        font.bold = is_bold
                        if shape == slide.shapes.title:
                            font.color.rgb = self.theme_color
                        else:
                            font.color.rgb = RGBColor(0, 0, 0)

TEMPLATES = {
    "Classic Blue": Template(
        "Classic Blue", RGBColor(0, 112, 192), RGBColor(255, 255, 255), 
        "Calibri", 40, 24, PP_ALIGN.LEFT, PP_ALIGN.LEFT, True, False
    ),
    "Modern Gradient": Template(
        "Modern Gradient", RGBColor(255, 255, 255), RGBColor(240, 240, 240), 
        "Arial", 44, 22, PP_ALIGN.CENTER, PP_ALIGN.LEFT, True, False,
        gradient=(RGBColor(0, 180, 180), RGBColor(0, 50, 100))
    ),
    "Minimalist Black": Template(
        "Minimalist Black", RGBColor(255, 255, 255), RGBColor(0, 0, 0), 
        "Helvetica", 48, 20, PP_ALIGN.LEFT, PP_ALIGN.LEFT, False, False
    ),
    "Vibrant Orange": Template(
        "Vibrant Orange", RGBColor(51, 51, 51), RGBColor(255, 200, 100), 
        "Verdana", 46, 26, PP_ALIGN.CENTER, PP_ALIGN.LEFT, True, False
    ),
    "Tech Blue": Template(
        "Tech Blue", RGBColor(255, 255, 255), RGBColor(240, 248, 255), 
        "Consolas", 38, 22, PP_ALIGN.LEFT, PP_ALIGN.LEFT, True, False,
        gradient=(RGBColor(173, 216, 230), RGBColor(0, 0, 139))
    ),
    "Retro Red": Template(
        "Retro Red", RGBColor(255, 255, 255), RGBColor(255, 200, 200), 
        "Impact", 50, 28, PP_ALIGN.LEFT, PP_ALIGN.LEFT, True, False
    ),
    "Futuristic Silver": Template(
        "Futuristic Silver", RGBColor(0, 0, 0), RGBColor(230, 230, 230), 
        "Century Gothic", 42, 24, PP_ALIGN.RIGHT, PP_ALIGN.LEFT, False, False,
        gradient=(RGBColor(192, 192, 192), RGBColor(105, 105, 105))
    ),
    "Earthy Brown": Template(
        "Earthy Brown", RGBColor(255, 255, 255), RGBColor(210, 180, 140), 
        "Bookman Old Style", 40, 22, PP_ALIGN.LEFT, PP_ALIGN.LEFT, True, False,
        shape=MSO_SHAPE.TRAPEZOID
    )
}