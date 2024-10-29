from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

class Template:
    def __init__(self, name, theme_color, background_color, font_family, title_font_size, content_font_size, 
                 title_alignment=PP_ALIGN.LEFT, content_alignment=PP_ALIGN.LEFT, 
                 title_bold=False, content_bold=False, gradient=None, shape=None,
                 slide_master=None):
        self.name = name
        self.theme_color = theme_color  # RGB tuple, e.g., (0, 112, 192)
        self.background_color = background_color  # RGB tuple
        self.font_family = font_family
        self.title_font_size = title_font_size
        self.content_font_size = content_font_size
        self.title_alignment = title_alignment
        self.content_alignment = content_alignment
        self.title_bold = title_bold
        self.content_bold = content_bold
        self.gradient = gradient  # tuple of two RGB tuples
        self.shape = shape
        self.slide_master = slide_master  # For future use
    
    def apply_style(self, slide):
        background = slide.background
        fill = background.fill
        
        if self.gradient:
            fill.gradient()
            fill.gradient_stops[0].color.rgb = RGBColor(*self.gradient[0])
            fill.gradient_stops[1].color.rgb = RGBColor(*self.gradient[1])
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self.background_color)
    
        if self.shape:
            left = top = Inches(0)
            width = height = Inches(10)
            shape = slide.shapes.add_shape(self.shape, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self.theme_color)
            shape.line.color.rgb = RGBColor(*self.theme_color)
    
        for shp in slide.shapes:
            if hasattr(shp, 'text_frame'):
                for paragraph in shp.text_frame.paragraphs:
                    if shp == slide.shapes.title:
                        paragraph.alignment = self.title_alignment
                        font_size = self.title_font_size
                        is_bold = self.title_bold
                        bullet = False
                    else:
                        paragraph.alignment = self.content_alignment
                        font_size = self.content_font_size
                        is_bold = self.content_bold
                        bullet = True
                    for run in paragraph.runs:
                        font = run.font
                        font.name = self.font_family
                        font.size = Pt(font_size)
                        font.bold = is_bold
                        if shp == slide.shapes.title:
                            font.color.rgb = RGBColor(*self.theme_color)
                        else:
                            # Choose text color based on background
                            if self.background_color == (255, 255, 255):
                                font.color.rgb = RGBColor(0, 0, 0)  # Black text on white background
                            else:
                                font.color.rgb = RGBColor(255, 255, 255)  # White text on colored background
                        font.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Enable auto-sizing
                        paragraph.space_before = Pt(6)
                        paragraph.space_after = Pt(6)
                        paragraph.bullet = bullet

# Define improved templates with better color matching and minimal black text
TEMPLATES = {
    "Elegant Blue": Template(
        "Elegant Blue", (0, 102, 204), (255, 255, 255), 
        "Calibri", 44, 24, PP_ALIGN.CENTER, PP_ALIGN.LEFT, True, False,
        gradient=((0, 102, 204), (204, 229, 255))
    ),
    "Modern Gray": Template(
        "Modern Gray", (34, 34, 34), (255, 255, 255), 
        "Arial", 40, 22, PP_ALIGN.LEFT, PP_ALIGN.LEFT, True, False,
        gradient=((34, 34, 34), (102, 102, 102))
    ),
    "Vibrant Green": Template(
        "Vibrant Green", (0, 153, 76), (240, 255, 240), 
        "Verdana", 42, 20, PP_ALIGN.LEFT, PP_ALIGN.LEFT, True, False,
        gradient=((0, 153, 76), (144, 238, 144))
    ),
    "Creative Purple": Template(
        "Creative Purple", (102, 0, 204), (245, 245, 245), 
        "Helvetica", 46, 22, PP_ALIGN.CENTER, PP_ALIGN.LEFT, True, False,
        gradient=((102, 0, 204), (204, 153, 255))
    ),
    "Sleek Black": Template(
        "Sleek Black", (0, 0, 0), (245, 245, 245), 
        "Tahoma", 48, 24, PP_ALIGN.LEFT, PP_ALIGN.LEFT, False, False,
        shape=MSO_SHAPE.RECTANGLE
    ),
    "Sunny Yellow": Template(
        "Sunny Yellow", (255, 204, 0), (255, 255, 224), 
        "Century Gothic", 40, 20, PP_ALIGN.LEFT, PP_ALIGN.LEFT, True, False,
        gradient=((255, 204, 0), (255, 255, 224))
    ),
    "Ocean Teal": Template(
        "Ocean Teal", (0, 128, 128), (224, 255, 255), 
        "Calibri Light", 44, 22, PP_ALIGN.CENTER, PP_ALIGN.LEFT, True, False,
        gradient=((0, 128, 128), (175, 238, 238))
    ),
    "Minimalist White": Template(
        "Minimalist White", (255, 255, 255), (245, 245, 245), 
        "Segoe UI", 42, 20, PP_ALIGN.LEFT, PP_ALIGN.LEFT, False, False,
        shape=MSO_SHAPE.OVAL
    )
}
