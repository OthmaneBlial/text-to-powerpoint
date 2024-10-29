import re
import logging
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image
from pptx.enum.text import PP_ALIGN
import requests
import os

class SlideGenerator:
    def __init__(self, template):
        self.prs = Presentation()
        self.template = template
        self.layouts = {
            'title': self.prs.slide_layouts[0],
            'content': self.prs.slide_layouts[1],
            'section': self.prs.slide_layouts[2],
            'image': self.prs.slide_layouts[5],
            'chart': self.prs.slide_layouts[5],
            'quote': self.prs.slide_layouts[5],
        }

    def parse_input(self, input_text):
        slides = []
        lines = input_text.strip().split('\n')

        # Parse title and optional subtitle
        title_match = re.match(r'^# (.+)', lines[0]) if len(lines) > 0 else None
        subtitle_match = re.match(r'^## (.+)', lines[1]) if len(lines) > 1 else None
        if title_match:
            title = title_match.group(1).strip()
            subtitle = subtitle_match.group(1).strip() if subtitle_match else ""
            slides.append({"type": "title", "title": title, "subtitle": subtitle})
        else:
            pass

        current_slide = {"type": "", "title": "", "content": []}
        for line in lines[1:]:
            line = line.strip()
            if re.match(r'^# .+', line):
                if current_slide["title"]:
                    slides.append(current_slide)
                title = re.sub(r'^# ', '', line).strip()
                current_slide = {"type": "content", "title": title, "content": []}
            elif re.match(r'^## .+', line):
                if current_slide["title"]:
                    slides.append(current_slide)
                title = re.sub(r'^## ', '', line).strip()
                current_slide = {"type": "section", "title": title, "content": []}
            elif re.match(r'^!\[.*\]\((.*)\)', line):
                image_url = re.findall(r'^!\[.*\]\((.*)\)', line)[0]
                slides.append({"type": "image", "title": current_slide.get("title", "Image Slide"), "image": image_url})
            elif re.match(r'^@chart\s*\{(.+)\}', line):
                chart_data = re.findall(r'^@chart\s*\{(.+)\}', line)[0]
                slides.append({"type": "chart", "title": current_slide.get("title", "Chart Slide"), "chart_data": chart_data})
            elif re.match(r'^> .+', line):
                quote = re.sub(r'^> ', '', line).strip()
                slides.append({"type": "quote", "quote": quote})
            elif re.match(r'^- (.+)', line) or re.match(r'^\* (.+)', line):
                # Directly add the bullet point without adding an extra "."
                content = re.sub(r'^[-*] ', '', line).strip()
                current_slide["content"].append(content)
            elif line:
                current_slide["content"].append(line)
        if current_slide["title"]:
            slides.append(current_slide)
        return slides

    def create_slide(self, slide_data):
        slide_type = slide_data["type"]
        layout = self.layouts.get(slide_type, self.layouts['content'])
        slide = self.prs.slides.add_slide(layout)
        self.template.apply_style(slide)

        if slide_type == "title":
            slide.shapes.title.text = slide_data["title"]
            if "subtitle" in slide_data:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = slide_data["subtitle"]
        elif slide_type in ["content", "section"]:
            slide.shapes.title.text = slide_data["title"]
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for item in slide_data["content"]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(6)
                p.bullet = True
        elif slide_type == "image":
            slide.shapes.title.text = slide_data["title"]
            image_path = self.download_image(slide_data["image"])
            if image_path:
                left = Inches(1)
                top = Inches(1.5)
                height = Inches(5.5)
                slide.shapes.add_picture(image_path, left, top, height=height)
        elif slide_type == "chart":
            slide.shapes.title.text = slide_data["title"]
            chart_data = self.parse_chart_data(slide_data["chart_data"])
            if chart_data:
                x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                ).chart
        elif slide_type == "quote":
            slide.shapes.title.text = "Quote"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = slide_data["quote"]
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER

        return slide

    def download_image(self, url):
        try:
            response = requests.get(url)
            response.raise_for_status()
            image = Image.open(BytesIO(response.content))
            image_path = f"temp_image_{hash(url)}.png"
            image.save(image_path)
            return image_path
        except Exception as e:
            logging.error(f"Failed to download image from {url}: {e}")
            return None

    def parse_chart_data(self, data_str):
        try:
            # Expected format: Categories: A,B,C; Series1: 10,20,30; Series2: 15,25,35
            data = {}
            parts = data_str.split(';')
            for part in parts:
                key, values = part.split(':')
                data[key.strip()] = [float(v.strip()) for v in values.split(',')]
            categories = data.get("Categories", [])
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for series_name, values in data.items():
                if series_name != "Categories":
                    chart_data.add_series(series_name, values)
            return chart_data
        except Exception as e:
            logging.error(f"Failed to parse chart data: {e}")
            return None

    def generate_presentation(self, input_text, output_file):
        slides_data = self.parse_input(input_text)
        for slide_data in slides_data:
            self.create_slide(slide_data)
        self.prs.save(output_file)
